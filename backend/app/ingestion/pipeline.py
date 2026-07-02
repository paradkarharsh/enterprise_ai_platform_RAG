"""
Document ingestion pipeline: parsing, chunking, embedding, and graph extraction.
"""
import logging
import hashlib
import os
from typing import List, Optional, Dict, Any
from uuid import uuid4
from pydantic import BaseModel
from app.config import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()


class ParsedDocument(BaseModel):
    """Result of document parsing."""
    title: str
    content: str
    pages: List[str] = []
    metadata: Dict[str, Any] = {}
    page_count: int = 0


class TextChunk(BaseModel):
    """A chunk of text with metadata."""
    id: str
    content: str
    chunk_index: int
    page_number: Optional[int] = None
    start_char: int = 0
    end_char: int = 0
    token_count: int = 0
    metadata: Dict[str, Any] = {}


# ─────────────────────────────────────────────
# Parsers
# ─────────────────────────────────────────────

async def parse_pdf(file_path: str) -> ParsedDocument:
    """Parse PDF document."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return ParsedDocument(
            title=os.path.basename(file_path),
            content="\n\n".join(pages),
            pages=pages,
            page_count=len(pages),
            metadata={"source_type": "pdf", "file_path": file_path},
        )
    except Exception as e:
        logger.error("PDF parsing failed: %s", e)
        raise


async def parse_docx(file_path: str) -> ParsedDocument:
    """Parse DOCX document."""
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return ParsedDocument(
            title=os.path.basename(file_path),
            content="\n\n".join(paragraphs),
            pages=paragraphs,
            page_count=1,
            metadata={"source_type": "docx", "file_path": file_path},
        )
    except Exception as e:
        logger.error("DOCX parsing failed: %s", e)
        raise


async def parse_pptx(file_path: str) -> ParsedDocument:
    """Parse PPTX document."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slides.append("\n".join(texts))
        return ParsedDocument(
            title=os.path.basename(file_path),
            content="\n\n---\n\n".join(slides),
            pages=slides,
            page_count=len(slides),
            metadata={"source_type": "pptx", "file_path": file_path},
        )
    except Exception as e:
        logger.error("PPTX parsing failed: %s", e)
        raise


async def parse_txt(file_path: str) -> ParsedDocument:
    """Parse plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()
    return ParsedDocument(
        title=os.path.basename(file_path),
        content=content,
        pages=[content],
        page_count=1,
        metadata={"source_type": "txt", "file_path": file_path},
    )


async def parse_csv(file_path: str) -> ParsedDocument:
    """Parse CSV file."""
    import csv
    rows = []
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        headers = next(reader, [])
        for row in reader:
            row_text = ", ".join(f"{h}: {v}" for h, v in zip(headers, row) if v.strip())
            rows.append(row_text)
    return ParsedDocument(
        title=os.path.basename(file_path),
        content="\n".join(rows),
        pages=rows,
        page_count=1,
        metadata={"source_type": "csv", "file_path": file_path, "headers": headers},
    )


async def parse_excel(file_path: str) -> ParsedDocument:
    """Parse Excel file."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        all_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
            all_text.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))
        return ParsedDocument(
            title=os.path.basename(file_path),
            content="\n\n".join(all_text),
            pages=all_text,
            page_count=len(all_text),
            metadata={"source_type": "excel", "file_path": file_path},
        )
    except Exception as e:
        logger.error("Excel parsing failed: %s", e)
        raise


async def parse_website(url: str) -> ParsedDocument:
    """Parse website content."""
    try:
        import httpx
        from bs4 import BeautifulSoup
        async with httpx.AsyncClient() as client:
            response = await client.get(url, follow_redirects=True, timeout=30)
            soup = BeautifulSoup(response.text, "html.parser")
            # Remove scripts and styles
            for element in soup(["script", "style", "nav", "footer", "header"]):
                element.decompose()
            title = soup.title.string if soup.title else url
            text = soup.get_text(separator="\n", strip=True)
        return ParsedDocument(
            title=title,
            content=text,
            pages=[text],
            page_count=1,
            metadata={"source_type": "website", "url": url},
        )
    except Exception as e:
        logger.error("Website parsing failed: %s", e)
        raise


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_txt,
    ".csv": parse_csv,
    ".xlsx": parse_excel,
    ".xls": parse_excel,
}


async def parse_document(file_path: str) -> ParsedDocument:
    """Parse any supported document type."""
    ext = os.path.splitext(file_path)[1].lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")
    return await parser(file_path)


# ─────────────────────────────────────────────
# Chunking
# ─────────────────────────────────────────────

def chunk_text(
    text: str,
    chunk_size: int = None,
    chunk_overlap: int = None,
    metadata: Dict[str, Any] = None,
) -> List[TextChunk]:
    """Split text into overlapping chunks using recursive character splitting."""
    chunk_size = chunk_size or settings.CHUNK_SIZE
    chunk_overlap = chunk_overlap or settings.CHUNK_OVERLAP

    separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""]
    chunks = []

    def _split_text(text: str, separators: list) -> list:
        final_chunks = []
        separator = separators[-1]
        for sep in separators:
            if sep in text:
                separator = sep
                break

        splits = text.split(separator) if separator else list(text)
        current_chunk = []
        current_length = 0

        for split in splits:
            piece = split + separator if separator else split
            if current_length + len(piece) > chunk_size and current_chunk:
                chunk_text = "".join(current_chunk).strip()
                if chunk_text:
                    final_chunks.append(chunk_text)
                # Overlap: keep last portion
                overlap_text = "".join(current_chunk)
                current_chunk = []
                current_length = 0
                if chunk_overlap > 0:
                    overlap = overlap_text[-chunk_overlap:]
                    current_chunk = [overlap]
                    current_length = len(overlap)

            current_chunk.append(piece)
            current_length += len(piece)

        if current_chunk:
            chunk_text = "".join(current_chunk).strip()
            if chunk_text:
                final_chunks.append(chunk_text)

        return final_chunks

    raw_chunks = _split_text(text, separators)

    offset = 0
    for i, content in enumerate(raw_chunks):
        start = text.find(content, offset)
        if start == -1:
            start = offset
        end = start + len(content)

        chunks.append(TextChunk(
            id=str(uuid4()),
            content=content,
            chunk_index=i,
            start_char=start,
            end_char=end,
            token_count=len(content.split()),
            metadata=metadata or {},
        ))
        offset = end

    return chunks


def compute_file_hash(file_path: str) -> str:
    """Compute SHA256 hash of a file."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            sha256.update(block)
    return sha256.hexdigest()
